from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from io import BytesIO
from flask_cors import CORS
import requests
import os

app = Flask(__name__)
CORS(app, origins=["https://areaprompt.com"])

SHARED_SECRET = "slidegen-2024-key-Zx4r9Lp1"
TEMPLATE_DIR = "templates"

def load_template(style):
    filename = f"{style.lower()}.pptx"
    path = os.path.join(TEMPLATE_DIR, filename)
    if not os.path.exists(path):
        raise FileNotFoundError(f"Template non trovato: {filename}")
    return Presentation(path)

def _rgb(hex_color):
    if not hex_color or not isinstance(hex_color, str) or len(hex_color) < 6:
        return RGBColor(0, 0, 0)
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def remove_default_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

def convert_bullets(text):
    if not text:
        return []
    lines = text.split('\n')
    items = []
    for line in lines:
        line = line.strip()
        if line.startswith('- '):
            items.append(('li', line[2:].strip()))
        elif line:
            items.append(('p', line))
    return items

def get_layout_by_name(prs, name, fallback_index=0):
    for layout in prs.slide_layouts:
        if layout.name.lower().strip() == name.lower().strip():
            return layout
    print(f"[AVVISO] Layout '{name}' non trovato. Uso fallback slide_layouts[{fallback_index}].")
    return prs.slide_layouts[fallback_index]

def get_placeholder_by_type(slide, placeholder_type):
    for ph in slide.placeholders:
        if ph.placeholder_format.type == placeholder_type:
            return ph
    return None

def create_presentation(slides_data, title=None, style=None, format="16:9", dimensions=None, fonts=None):
    try:
        prs = load_template(style)
    except FileNotFoundError:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

    remove_default_slides(prs)

    for slide_info in slides_data:
        layout_name = slide_info.get("layout_name", "Titolo e contenuto")
        slide_layout = get_layout_by_name(prs, layout_name)
        slide = prs.slides.add_slide(slide_layout)

        # Inserisci titolo
        title_text = slide_info.get("title", "")
        placeholder_title = get_placeholder_by_type(slide, PP_PLACEHOLDER.TITLE)
        if placeholder_title:
            placeholder_title.text = title_text

        # Inserisci contenuto
        content_text = slide_info.get("content", "")
        placeholder_content = get_placeholder_by_type(slide, PP_PLACEHOLDER.BODY)
        if placeholder_content:
            content_frame = placeholder_content.text_frame
            content_frame.clear()
            for type_, txt in convert_bullets(content_text):
                para = content_frame.add_paragraph()
                para.text = txt
                if type_ == 'li':
                    para.level = 0

        # Inserisci immagine opzionale
        image_url = slide_info.get("image_url")
        if image_url:
            try:
                response = requests.get(image_url, timeout=10)
                response.raise_for_status()
                image_stream = BytesIO(response.content)

                # Posizione immagine dinamica
                layout_pos = {
                    "Immagine destra + Testo sinistra": (Inches(6), Inches(1.5)),
                    "Immagine sinistra + Testo destra": (Inches(0.5), Inches(1.5)),
                    "Immagine centrata + Testo sotto": (Inches(3.5), Inches(1.0))
                }
                left, top = layout_pos.get(layout_name, (Inches(6), Inches(1.5)))

                slide.shapes.add_picture(image_stream, left, top, Inches(3.5), Inches(4.0))
            except Exception as e:
                print(f"Image error: {e}")
                placeholder = slide.shapes.add_textbox(Inches(6), Inches(1.5), Inches(3.5), Inches(1))
                placeholder.text_frame.text = "Image not available"

        # Fallback se mancano tutti i placeholder
        if not placeholder_title and not placeholder_content:
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
            tf = textbox.text_frame
            tf.text = f"{title_text}\n\n{content_text}"

    return prs

@app.route("/generate", methods=["POST"])
def generate_pptx():
    data = request.get_json()
    if not data or "slides" not in data or data.get("secret") != SHARED_SECRET:
        return jsonify({"error": "Unauthorized or invalid input"}), 403

    try:
        prs = create_presentation(
            slides_data=data["slides"],
            title=data.get("title"),
            style=data.get("style"),
            format=data.get("format"),
            dimensions=data.get("dimensions"),
            fonts=data.get("fonts") or {}
        )

        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        return send_file(
            pptx_io,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            as_attachment=True,
            download_name=f"presentazione_{data.get('style','default')}.pptx"
        )
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    # Debug: stampa i layout del primo template disponibile
    try:
        test_prs = load_template("default")
        print("Layout disponibili:")
        for i, layout in enumerate(test_prs.slide_layouts):
            print(f"{i}: {layout.name}")
    except:
        pass

    app.run(debug=True)

